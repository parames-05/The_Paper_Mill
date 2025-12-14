from datetime import datetime, timezone
import os
import tempfile
import shutil
import io
import pikepdf
import traceback
from flask import Flask, render_template, request, send_file, abort, url_for
import fitz
from PIL import Image
import pillow_heif
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from reportlab.pdfgen import canvas
import file_handler
import pandas as pd
import csv
import chardet

# This is CRUCIAL for HEIF/HEIC support
try:
    from pillow_heif import register_heif_opener

    register_heif_opener()
except ImportError:
    print("WARNING: pillow-heif not installed. HEIF/HEIC conversion will fail.")
    pass

# Import our custom file handler
import file_handler

app = Flask(__name__)
pillow_heif.register_heif_opener()



# --- Page Routes ---
# ... (all your @app.route page routes are correct) ...

@app.route("/healthz")
def health_check():
    return {
        "status": "ok",
        "service": "The_Paper_Mill",
        "time": datetime.now(timezone.utc).isoformat()
    }, 200


@app.route("/work-in-progress")
def work_in_progress():
    return render_template("work_in_progress.html")

@app.route("/how-to-use")
def how_to_use():
    return render_template("how_to_use.html")

@app.route('/')
def index():
    # Redirect to the first converter page
    return render_template("index.html")


@app.route('/jpg-to-pdf')
def jpg_to_pdf():
    return render_template('converter_page.html',
                           title='JPG to PDF',
                           file_accept='.jpg')


@app.route('/jpeg-to-pdf')
def jpeg_to_pdf():
    return render_template('converter_page.html',
                           title='JPEG to PDF',
                           file_accept='.jpeg, .jpg')


@app.route('/png-to-pdf')
def png_to_pdf():
    return render_template('converter_page.html',
                           title='PNG to PDF',
                           file_accept='.png')


@app.route('/bmp-to-pdf')
def bmp_to_pdf():
    return render_template('converter_page.html',
                           title='BMP to PDF',
                           file_accept='.bmp')


@app.route('/tiff-to-pdf')
def tiff_to_pdf():
    return render_template('converter_page.html',
                           title='TIFF to PDF',
                           file_accept='.tiff, .tif')


@app.route('/webp-to-pdf')
def webp_to_pdf():
    return render_template('converter_page.html',
                           title='WebP to PDF',
                           file_accept='.webp')


@app.route('/heif-to-pdf')
def heif_to_pdf():
    return render_template('converter_page.html',
                           title='HEIF to PDF',
                           file_accept='.heif, .heic')


# --- Core Conversion Route ---
# All forms will post to this single endpoint

@app.route('/convert-images', methods=['POST'])
def convert_images_to_pdf():
    # 1. Use our modular handler to save files
    temp_dir = None
    try:
        temp_dir, image_paths = file_handler.save_uploaded_files(request)

        if not image_paths:
            return "No images were uploaded.", 400

        # 2. Process the images
        pil_images = []
        for path in image_paths:
            try:
                img = Image.open(path)

                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')

                pil_images.append(img)
            except Exception as e:
                print(f"Error opening image {path}: {e}")
                pass

        if not pil_images:
            return "Could not read any of the uploaded images.", 400

        # --- START OF FIX ---

        # 3. Save as PDF... in memory!
        pdf_buffer = io.BytesIO()

        # Save the first image, and append the rest
        pil_images[0].save(
            pdf_buffer,
            format='PDF',  # Specify format when saving to buffer
            save_all=True,
            append_images=pil_images[1:]
        )

        # Rewind the buffer to the beginning so send_file can read it
        pdf_buffer.seek(0)

        # 4. Send the PDF buffer to the user
        return send_file(
            pdf_buffer,
            as_attachment=False,  # <-- Changed from True to False
            download_name='convertigo.pdf',  # This is still good to have
            mimetype='application/pdf'  # Specify mimetype
        )

        # --- END OF FIX ---

    except Exception as e:
        print(f"An error occurred during conversion: {e}")
        abort(500, description="An internal error occurred during conversion.")

    finally:
        # 5. Clean up the temporary directory (for the uploaded images)
        # This is now safe, as the PDF was never in this directory.
        if temp_dir and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)


# ==========================
# PDF MERGE
# ==========================

@app.route('/merge-pdf')
def merge_pdf():
    return render_template(
        'converter_page.html',
        title='Merge PDFs',
        file_accept='.pdf',
        upload_endpoint=url_for('merge_pdfs')
    )


@app.route('/merge', methods=['POST'])
def merge_pdfs():
    temp_dir, pdf_paths = file_handler.save_uploaded_files(request)
    merger = PdfMerger()
    try:
        for p in pdf_paths:
            merger.append(p)
        output = io.BytesIO()
        merger.write(output)
        output.seek(0)
        return send_file(output, as_attachment=True, download_name='merged.pdf', mimetype='application/pdf')
    except Exception as e:
        print(f"Merge error: {e}")
        abort(500)
    finally:
        merger.close()
        shutil.rmtree(temp_dir)

# ==========================
# PDF SPLIT
# ==========================

@app.route('/split-pdf')
def split_pdf():
    return render_template(
        'converter_page.html',
        title='Split PDF',
        file_accept='.pdf',
        upload_endpoint=url_for('split_pdf_action')
    )


@app.route('/split', methods=['POST'])
def split_pdf_action():
    temp_dir, pdf_paths = file_handler.save_uploaded_files(request)
    if not pdf_paths:
        return "No PDF uploaded.", 400

    input_pdf = PdfReader(pdf_paths[0])
    split_dir = tempfile.mkdtemp()

    for i, page in enumerate(input_pdf.pages):
        writer = PdfWriter()
        writer.add_page(page)
        with open(os.path.join(split_dir, f"page_{i+1}.pdf"), "wb") as f:
            writer.write(f)

    zip_path = shutil.make_archive(os.path.join(split_dir, "split_pages"), 'zip', split_dir)
    return send_file(zip_path, as_attachment=True, download_name="split_pages.zip")

# ==========================
# WORD → PDF
# ==========================

@app.route('/word-to-pdf')
def word_to_pdf():
    return render_template(
        'converter_page.html',
        title='Word to PDF',
        file_accept='.docx',
        upload_endpoint=url_for('convert_word_to_pdf'),
        show_doc_warning=True
    )

@app.route('/convert-word', methods=['POST'])
def convert_word_to_pdf():
    import mammoth
    import io
    import os
    import shutil
    from flask import send_file, after_this_request, request
    from playwright.sync_api import sync_playwright

    # Save uploaded file
    temp_dir, doc_paths = file_handler.save_uploaded_files(request)
    if not doc_paths:
        return "No DOCX uploaded", 400

    doc_path = doc_paths[0]
    html_path = os.path.join(temp_dir, "document.html")

    # Cleanup AFTER response is fully sent
    @after_this_request
    def cleanup(response):
        try:
            shutil.rmtree(temp_dir)
        except Exception as e:
            print("Cleanup warning:", e)
        return response

    # ==========================
    # DOCX → HTML
    # ==========================
    with open(doc_path, "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file)

        html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <style>
                @page {{
                    size: A4;
                    margin: 40px;
                }}

                body {{
                    font-family: Arial, Helvetica, sans-serif;
                    font-size: 12pt;
                    line-height: 1.4;
                }}

                img {{
                    max-width: 100%;
                }}

                table {{
                    width: 100%;
                    border-collapse: collapse;
                }}

                table, th, td {{
                    border: 1px solid #444;
                }}

                th, td {{
                    padding: 6px;
                }}

                /* Respect Word-style page breaks if present */
                .page-break {{
                    page-break-before: always;
                }}
            </style>
        </head>
        <body>
            {result.value}
        </body>
        </html>
        """

        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html)

    # ==========================
    # HTML → PDF (Chromium Print Engine)
    # ==========================
    pdf_buffer = io.BytesIO()

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page()

        page.goto(f"file:///{html_path}")
        page.emulate_media(media="print")

        pdf_bytes = page.pdf(
            format="A4",
            print_background=True,
            margin={
                "top": "40px",
                "bottom": "40px",
                "left": "40px",
                "right": "40px"
            }
        )

        browser.close()

    pdf_buffer.write(pdf_bytes)
    pdf_buffer.seek(0)

    # ==========================
    # Send PDF
    # ==========================
    return send_file(
        pdf_buffer,
        as_attachment=True,
        download_name="word_to_pdf.pdf",
        mimetype="application/pdf"
    )

# ==========================
# EXCEL → PDF
# ==========================

@app.route('/excel-to-pdf')
def excel_to_pdf():
    return render_template(
        'converter_page.html',
        title='Excel to PDF',
        file_accept='.xlsx',
        upload_endpoint=url_for('convert_excel_to_pdf'),
        show_doc_warning=True
    )


@app.route('/convert-excel', methods=['POST'])
def convert_excel_to_pdf():
    from openpyxl import load_workbook
    temp_dir, excel_paths = file_handler.save_uploaded_files(request)
    excel_path = excel_paths[0]

    wb = load_workbook(excel_path)
    sheet = wb.active
    pdf_buffer = io.BytesIO()
    c = canvas.Canvas(pdf_buffer)
    y = 800

    for row in sheet.iter_rows(values_only=True):
        row_text = " | ".join([str(cell) if cell else "" for cell in row])
        c.drawString(50, y, row_text)
        y -= 15
        if y < 50:
            c.showPage()
            y = 800

    c.save()
    pdf_buffer.seek(0)
    return send_file(pdf_buffer, as_attachment=True, download_name='excel_to_pdf.pdf', mimetype='application/pdf')

# ==========================
# POWERPOINT → PDF
# ==========================

@app.route('/pptx-to-pdf')
def pptx_to_pdf():
    return render_template(
        'converter_page.html',
        title='PPTX to PDF',
        file_accept='.pptx',
        upload_endpoint=url_for('convert_pptx_to_pdf'),
        show_doc_warning=True
    )

@app.route('/convert-pptx', methods=['POST'])
def convert_pptx_to_pdf():
    import io
    import os
    import shutil
    import base64
    from flask import send_file, after_this_request, request
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from playwright.sync_api import sync_playwright

    temp_dir, ppt_paths = file_handler.save_uploaded_files(request)
    if not ppt_paths:
        return "No PPTX uploaded", 400

    ppt_path = ppt_paths[0]
    html_path = os.path.join(temp_dir, "slides.html")
    images_dir = os.path.join(temp_dir, "images")
    os.makedirs(images_dir, exist_ok=True)

    @after_this_request
    def cleanup(response):
        try:
            shutil.rmtree(temp_dir)
        except Exception as e:
            print("Cleanup warning:", e)
        return response

    prs = Presentation(ppt_path)
    slide_blocks = []

    # ==========================
    # SLIDE PARSING
    # ==========================
    for slide_idx, slide in enumerate(prs.slides):
        elements = []

        for shape in slide.shapes:

            # ---------- TEXT ----------
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    elements.append(f"<p class='text'>{text}</p>")

            # ---------- IMAGES ----------
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_name = f"slide{slide_idx}_{shape.shape_id}.{image_ext}"
                image_path = os.path.join(images_dir, image_name)

                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)

                elements.append(
                    f"<img src='images/{image_name}' class='slide-image' />"
                )

            # ---------- CHARTS / SMARTART / EVERYTHING ELSE ----------
            else:
                # Ignore unsupported shapes safely
                continue

        slide_html = f"""
        <section class="slide">
            {''.join(elements)}
        </section>
        """
        slide_blocks.append(slide_html)

    # ==========================
    # HTML TEMPLATE
    # ==========================
    html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="utf-8">
        <style>
            @page {{
                size: 16:9;
                margin: 0;
            }}

            body {{
                margin: 0;
                font-family: Arial, Helvetica, sans-serif;
                background: white;
            }}

            .slide {{
                width: 960px;
                height: 540px;
                padding: 40px;
                box-sizing: border-box;
                page-break-after: always;
                display: flex;
                flex-direction: column;
                gap: 12px;
            }}

            .text {{
                font-size: 22px;
            }}

            .slide-image {{
                max-width: 100%;
                max-height: 320px;
                object-fit: contain;
            }}

            .unsupported {{
                margin-top: 20px;
                padding: 10px;
                border: 1px dashed red;
                color: red;
                font-size: 14px;
            }}
        </style>
    </head>
    <body>
        {''.join(slide_blocks)}
    </body>
    </html>
    """

    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html)

    # ==========================
    # HTML → PDF
    # ==========================
    pdf_buffer = io.BytesIO()

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page()
        page.goto(f"file:///{html_path}")
        page.emulate_media(media="print")

        pdf_bytes = page.pdf(
            width="960px",
            height="540px",
            print_background=True
        )

        browser.close()

    pdf_buffer.write(pdf_bytes)
    pdf_buffer.seek(0)

    return send_file(
        pdf_buffer,
        as_attachment=True,
        download_name="pptx_to_pdf.pdf",
        mimetype="application/pdf"
    )


# ==========================
# HEIC → PDF
# ==========================

@app.route('/heic-to-pdf')
def heic_to_pdf_page():
    return render_template(
        'converter_page.html',
        title='HEIC to PDF',
        file_accept='.heic,.heif',
        upload_endpoint=url_for('convert_heic_to_pdf')
    )

@app.route('/convert-heic-to-pdf', methods=['POST'])
def convert_heic_to_pdf():
    temp_dir, heic_paths = file_handler.save_uploaded_files(request)
    images = []

    try:
        for path in heic_paths:
            img = Image.open(path)
            img = img.convert("RGB")  # ensure proper PDF compatibility
            images.append(img)

        pdf_path = os.path.join(temp_dir, "converted_heic.pdf")
        images[0].save(pdf_path, save_all=True, append_images=images[1:])
        return send_file(pdf_path, as_attachment=True, download_name="converted_heic.pdf")
    finally:
        shutil.rmtree(temp_dir)

# ==========================
# ZIP (FOLDER OF IMAGES) → PDF
# ==========================

@app.route('/zip-to-pdf')
def zip_to_pdf():
    return render_template(
        'converter_page.html',
        title='ZIP (Images) to PDF',
        file_accept='.zip',
        upload_endpoint=url_for('convert_zip_to_pdf')
    )

@app.route('/convert-zip-to-pdf', methods=['POST'])
def convert_zip_to_pdf():
    import zipfile
    temp_dir, zip_paths = file_handler.save_uploaded_files(request)
    if not zip_paths:
        return "No ZIP file uploaded.", 400

    zip_path = zip_paths[0]
    extract_dir = tempfile.mkdtemp()

    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(extract_dir)

        image_files = []
        for root, _, files in os.walk(extract_dir):
            for f in files:
                if file_handler.allowed_file(f):
                    path = os.path.join(root, f)
                    if os.path.splitext(f.lower())[1] in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp']:
                        image_files.append(path)

        if not image_files:
            return "No supported images found inside the ZIP.", 400

        image_files.sort()
        pil_images = []
        for path in image_files:
            img = Image.open(path)
            if img.mode in ('RGBA', 'P'):
                img = img.convert('RGB')
            pil_images.append(img)

        pdf_buffer = io.BytesIO()
        pil_images[0].save(pdf_buffer, format='PDF', save_all=True, append_images=pil_images[1:])
        pdf_buffer.seek(0)
        return send_file(pdf_buffer, as_attachment=True, download_name='zip_to_pdf.pdf', mimetype='application/pdf')
    finally:
        shutil.rmtree(temp_dir)
        shutil.rmtree(extract_dir)


# ==========================
# PDF COMPRESSOR (with pikepdf)
# ==========================

@app.route('/compress-pdf')
def compress_pdf_page():
    return render_template(
        'converter_page.html',
        title='Compress PDF',
        file_accept='.pdf',
        upload_endpoint=url_for('compress_pdf_action')
    )


@app.route('/compress-pdf-action', methods=['POST'])
def compress_pdf_action():
    temp_dir, pdf_paths = file_handler.save_uploaded_files(request)
    if not pdf_paths:
        shutil.rmtree(temp_dir)
        return "No valid PDF uploaded", 400

    input_path = pdf_paths[0]
    output_path = os.path.join(temp_dir, "compressed.pdf")

    try:
        # Open original PDF
        doc = fitz.open(input_path)
        new_pdf = fitz.open()

        # Render each page as an image at lower DPI
        for page in doc:
            pix = page.get_pixmap(matrix=fitz.Matrix(1.2, 1.2))  # 1.5x reduces size a lot but keeps quality
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img_io = io.BytesIO()
            img.save(img_io, format="JPEG", optimize=True, quality=60)  # lower quality => smaller size
            img_io.seek(0)

            # Add new page to new PDF
            img_pdf = fitz.open(stream=img_io.read(), filetype="jpeg")
            rect = fitz.Rect(0, 0, pix.width, pix.height)
            page_new = new_pdf.new_page(width=rect.width, height=rect.height)
            page_new.insert_image(rect, stream=img_io.getvalue())

        new_pdf.save(output_path, garbage=4, deflate=True)
        new_pdf.close()
        doc.close()

        return send_file(
            output_path,
            as_attachment=True,
            download_name="compressed.pdf",
            mimetype="application/pdf"
        )

    except Exception as e:
        traceback.print_exc()
        return f"Compression failed: {e}", 500

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
# =========================
# IMAGE COMPRESSOR
# ==========================

@app.route('/compress-image')
def compress_image_page():
    return render_template(
        'converter_page.html',
        title='Compress Image',
        file_accept='.jpg,.jpeg,.png,.webp',
        upload_endpoint=url_for('compress_image_action')
    )


@app.route('/compress-image-action', methods=['POST'])
def compress_image_action():
    temp_dir, img_paths = file_handler.save_uploaded_files(request)
    if not img_paths:
        return "No image uploaded.", 400

    img_path = img_paths[0]
    quality = request.form.get('quality', default=60, type=int)

    # --- Solution: Use an in-memory buffer ---
    img_buffer = io.BytesIO()

    try:
        img = Image.open(img_path)

        # Handle transparency before saving as JPG
        if img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')

        # Save to the buffer, not a file
        img.save(img_buffer, format='JPEG', optimize=True, quality=quality)

        # IMPORTANT: Rewind the buffer to the beginning
        img_buffer.seek(0)

        # Send the buffer
        return send_file(
            img_buffer,
            as_attachment=True,
            download_name='compressed_image.jpg',
            mimetype='image/jpeg'  # It's good to be explicit
        )
    finally:
        # This now ONLY cleans up the original uploaded file's directory.
        # The compressed file was never saved to disk, so no lock!
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)



# ==========================
# CSV → XLSX
# ==========================

@app.route('/csv-to-xlsx')
def csv_to_xlsx_page():
    return render_template(
        'converter_page.html',
        title='CSV to XLSX',
        file_accept='.csv',
        upload_endpoint=url_for('convert_csv_to_xlsx'),
        show_doc_warning=True
    )

@app.route('/convert-csv-to-xlsx', methods=['POST'])
def convert_csv_to_xlsx():
    import pandas as pd
    import csv
    import chardet
    import os, shutil

    temp_dir, csv_paths = file_handler.save_uploaded_files(request)
    if not csv_paths:
        shutil.rmtree(temp_dir)
        return "No CSV file uploaded", 400

    csv_path = csv_paths[0]
    xlsx_path = os.path.join(temp_dir, "converted.xlsx")

    try:
        # Detect encoding
        with open(csv_path, 'rb') as f:
            raw = f.read(10000)
            encoding = chardet.detect(raw)['encoding'] or 'utf-8'

        # Read file and sniff delimiter
        with open(csv_path, 'r', encoding=encoding, errors='replace') as f:
            sample = f.read(2048)
            f.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=[',', ';', '\t', '|'])
                delimiter = dialect.delimiter
            except Exception:
                delimiter = ','

            df = pd.read_csv(f, delimiter=delimiter, encoding=encoding)

        # Save to Excel
        df.to_excel(xlsx_path, index=False, engine='openpyxl')

        return send_file(
            xlsx_path,
            as_attachment=True,
            download_name="converted.xlsx",
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    except Exception as e:
        return f"Conversion failed: {e}", 500

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


# ==========================
# JSON → CSV
# ==========================

@app.route('/json-to-csv')
def json_to_csv():
    return render_template(
        'converter_page.html',
        title='JSON to CSV',
        file_accept='.json',
        upload_endpoint=url_for('convert_json_to_csv'),
        show_doc_warning=True
    )

@app.route('/convert-json-to-csv', methods=['POST'])
def convert_json_to_csv():
    import pandas as pd
    import json
    temp_dir, json_paths = file_handler.save_uploaded_files(request)
    json_path = json_paths[0]

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Normalize to handle nested structures
    df = pd.json_normalize(data)
    csv_path = os.path.join(temp_dir, "converted.csv")
    df.to_csv(csv_path, index=False)
    return send_file(csv_path, as_attachment=True, download_name="converted.csv")


# ==========================
# MAIN
# ==========================

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=10000)
